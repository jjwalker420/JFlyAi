#!/usr/bin/env python3
"""Builds The Vision deck for Mark's launch packet.

Clean, modern, lots of white space, one idea per slide.
Output: The-Vision.pptx (converted to PDF separately via LibreOffice).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Palette
INK    = RGBColor(0x14, 0x30, 0x4A)  # deep navy, primary text
ACCENT = RGBColor(0xC8, 0x77, 0x2E)  # warm copper
CREAM  = RGBColor(0xFA, 0xF7, 0xF2)  # warm off white background
SAND   = RGBColor(0xEC, 0xE3, 0xD6)  # soft panel
SLATE  = RGBColor(0x5B, 0x6B, 0x7A)  # secondary text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DEEP   = RGBColor(0x0E, 0x24, 0x39)  # darker navy for title bg

FONT = "Liberation Sans"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s


def rect(s, x, y, w, h, color, line=None, line_w=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w or 1)
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.name = FONT
            r.font.color.rgb = color; r.font.bold = bold
    return tb


def kicker(s, label, num):
    rect(s, 0.9, 0.85, 0.45, 0.12, ACCENT)
    txt(s, 0.9, 1.0, 9.0, 0.4, [[(label.upper(), 13, ACCENT, True)]])
    # page number
    txt(s, SW - 1.5, SH - 0.8, 0.9, 0.4, [[(num, 12, SLATE, False)]], align=PP_ALIGN.RIGHT)
    txt(s, 0.9, SH - 0.8, 6.0, 0.4, [[("Your Documents Business", 11, SLATE, False)]])


def headline(s, lines, y=1.45, size=40, color=INK):
    txt(s, 0.9, y, 11.5, 1.6, [[(l, size, color, True)] for l in lines], line_spacing=1.02, space_after=2)


# ----------------------------------------------------------------------------
# Slide 1: Title
# ----------------------------------------------------------------------------
s = slide(DEEP)
rect(s, 0, 0, SW, SH, DEEP)
rect(s, 0.9, 2.35, 0.9, 0.16, ACCENT)
txt(s, 0.9, 1.2, 11, 0.5, [[("A BUSINESS BUILT FOR YOU", 15, ACCENT, True)]])
txt(s, 0.9, 2.7, 11.5, 2.2, [
    [("Turn what you already", 50, WHITE, True)],
    [("know into work you can", 50, WHITE, True)],
    [("do from home.", 50, ACCENT, True)],
], line_spacing=1.02, space_after=2)
txt(s, 0.9, 5.7, 11, 0.6, [[("Prepared for Mark Strazisar", 18, WHITE, False)]])
txt(s, 0.9, 6.25, 11, 0.6, [[("25 years of hospitality, finally on your terms.", 15, SLATE, False)]])

# ----------------------------------------------------------------------------
# Slide 2: The business in one sentence
# ----------------------------------------------------------------------------
s = slide()
kicker(s, "The idea", "2")
headline(s, ["The business, in one sentence."])
rect(s, 0.9, 3.0, 11.5, 2.4, WHITE)
rect(s, 0.9, 3.0, 0.14, 2.4, ACCENT)
txt(s, 1.45, 3.35, 10.6, 1.9, [
    [("You help busy restaurant groups turn what lives in their", 25, INK, False)],
    [("best people's heads into clear written playbooks, so every", 25, INK, False)],
    [("location runs the same and every new opening starts from a", 25, INK, False)],
    [("proven plan instead of a blank page.", 25, INK, True)],
], line_spacing=1.1, space_after=2)
txt(s, 0.9, 5.95, 11, 0.6, [[("No cooking. No floor. No clock. Just your knowledge, written down.", 15, SLATE, False)]])

# ----------------------------------------------------------------------------
# Slide 3: What your day looks like
# ----------------------------------------------------------------------------
s = slide()
kicker(s, "Your day", "3")
headline(s, ["What your day actually looks like."])
steps = [
    ("1", "Talk", "Open Claude and talk for about ten minutes. Explain how a task really gets done, the way you would tell a new manager."),
    ("2", "Review", "Claude writes a clean draft while you talk. You read it and fix anything that is not how you really do it."),
    ("3", "Save", "When it sounds like you, save it. That is one finished, sellable document. You are done for the day if you want."),
]
cw, gap, x0, y0, ch = 3.55, 0.35, 0.9, 3.0, 3.0
for i, (n, t, d) in enumerate(steps):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, WHITE)
    rect(s, x, y0, cw, 0.12, ACCENT)
    txt(s, x + 0.35, y0 + 0.4, cw - 0.7, 0.9, [[(n, 40, SAND, True)]])
    txt(s, x + 0.35, y0 + 1.0, cw - 0.7, 0.6, [[(t, 23, INK, True)]])
    txt(s, x + 0.35, y0 + 1.65, cw - 0.7, 1.2, [[(d, 13.5, SLATE, False)]], line_spacing=1.12)
txt(s, 0.9, 6.35, 11, 0.5, [[("That is the whole rhythm. Talk, review, save. You stay in charge of every word.", 14, INK, False)]])

# ----------------------------------------------------------------------------
# Slide 4: Why you
# ----------------------------------------------------------------------------
s = slide()
kicker(s, "Why you", "4")
headline(s, ["You have already built", "the thing you will sell."])
pts = [
    "Around 15 concepts built and operated: restaurants, bars, lounges, nightclubs, a coffee shop, a family fun center.",
    "Owner, Chief Operating Officer, and Vice President of Operations.",
    "You wrote the training manuals, the safety procedures, and the onboarding programs yourself.",
    "You took Beta Nightclub to number one in North America and doubled sales at SUTRA Lounge.",
    "You ran the P&Ls and mentored the general managers who ran your floors.",
]
y = 3.25
for p in pts:
    rect(s, 0.95, y + 0.07, 0.22, 0.22, ACCENT)
    txt(s, 1.45, y, 11.0, 0.7, [[(p, 16.5, INK, False)]], line_spacing=1.08)
    y += 0.75
txt(s, 0.9, 6.95, 11.5, 0.5, [[("This is not a new career. It is the work you have always been great at.", 14, SLATE, False)]])

# ----------------------------------------------------------------------------
# Slide 5: The problem your customer has
# ----------------------------------------------------------------------------
s = slide()
kicker(s, "The opportunity", "5")
headline(s, ["The problem every growing", "group has, and you solve."])
probs = [
    ("Training drifts", "A server learns one way at one spot and a different way at another. The brand feels different depending on the door you walk in."),
    ("Knowledge walks out", "The real playbook lives in your best people's heads. When they leave, it leaves with them."),
    ("Every opening from scratch", "A group that keeps opening rebuilds training and checklists each time, instead of pulling a proven playbook off the shelf."),
]
cw, gap, x0, y0, ch = 3.55, 0.35, 0.9, 3.2, 2.9
for i, (t, d) in enumerate(probs):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, WHITE)
    rect(s, x, y0, 0.12, ch, ACCENT)
    txt(s, x + 0.4, y0 + 0.4, cw - 0.75, 0.9, [[(t, 19, INK, True)]], line_spacing=1.0)
    txt(s, x + 0.4, y0 + 1.35, cw - 0.75, 1.4, [[(d, 13.5, SLATE, False)]], line_spacing=1.14)
txt(s, 0.9, 6.45, 11.5, 0.5, [[("Your lane is floor operations. You make the knowledge into documents anyone can use.", 14, INK, False)]])

# ----------------------------------------------------------------------------
# Slide 6: What you sell
# ----------------------------------------------------------------------------
s = slide()
kicker(s, "The offer", "6")
headline(s, ["What you sell."])
fams = [
    ("Operations core", "Opening and closing, cash handling, cleaning, inventory, food safety, emergency procedures, manager checklists."),
    ("Employee and training", "Onboarding, server and bartender and host manuals, steps of service, menu knowledge, coaching."),
    ("Standards and scaling", "Brand standards, audits, recipe specs, and the flagship: the new location opening playbook."),
]
cw, gap, x0, y0, ch = 3.55, 0.35, 0.9, 2.9, 2.9
for i, (t, d) in enumerate(fams):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, WHITE)
    rect(s, x, y0, cw, 0.12, ACCENT)
    txt(s, x + 0.38, y0 + 0.45, cw - 0.7, 0.9, [[(t, 18.5, INK, True)]], line_spacing=1.0)
    txt(s, x + 0.38, y0 + 1.35, cw - 0.7, 1.4, [[(d, 13.5, SLATE, False)]], line_spacing=1.14)
txt(s, 0.9, 6.2, 11.5, 0.7, [
    [("Sell one document, a package, or a full system. Sell it once, or keep it current for a monthly fee.", 14, INK, False)]])

# ----------------------------------------------------------------------------
# Slide 7: The flagship
# ----------------------------------------------------------------------------
s = slide(DEEP)
rect(s, 0, 0, SW, SH, DEEP)
rect(s, 0.9, 0.85, 0.45, 0.12, ACCENT)
txt(s, 0.9, 1.0, 9, 0.4, [[("YOUR FLAGSHIP", 13, ACCENT, True)]])
txt(s, SW - 1.5, SH - 0.8, 0.9, 0.4, [[("7", 12, SLATE, False)]], align=PP_ALIGN.RIGHT)
txt(s, 0.9, 1.5, 11.5, 1.6, [
    [("The new location", 40, WHITE, True)],
    [("opening playbook.", 40, ACCENT, True)],
], line_spacing=1.02, space_after=2)
txt(s, 0.9, 3.2, 11.3, 1.0, [
    [("Most groups rebuild every opening from a blank page. You give them one", 17, WHITE, False)],
    [("proven plan they run every time. A countdown timeline, build sheets, the", 17, WHITE, False)],
    [("hiring and training calendar, punch lists, and the soft opening plan.", 17, WHITE, False)],
], line_spacing=1.18)
rect(s, 0.9, 5.0, 11.5, 1.4, RGBColor(0x1B, 0x3B, 0x57))
rect(s, 0.9, 5.0, 0.14, 1.4, ACCENT)
txt(s, 1.4, 5.3, 10.8, 0.9, [
    [("It turns every opening from a scramble into a routine. ", 17, WHITE, False),
     ("Lead with this.", 17, ACCENT, True)],
], anchor=MSO_ANCHOR.MIDDLE)

# ----------------------------------------------------------------------------
# Slide 8: The money
# ----------------------------------------------------------------------------
s = slide()
kicker(s, "The money", "8")
headline(s, ["What the money looks like."])
ways = [
    ("Per document", "$150 to $750", "Each finished document. Great for starting small.", False),
    ("Monthly retainer", "$3,000 / month", "Recommended. A flat fee for the build period, then $1,000 a month to keep it current. Steady income, protected on slow weeks.", True),
    ("Project fee", "$18k to $25k", "One fixed price for the full system, paid in milestones.", False),
]
cw, gap, x0, y0, ch = 3.55, 0.35, 0.9, 2.85, 2.9
for i, (t, price, d, rec) in enumerate(ways):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, WHITE if not rec else INK)
    rect(s, x, y0, cw, 0.12, ACCENT)
    tcol = INK if not rec else WHITE
    scol = SLATE if not rec else SAND
    if rec:
        txt(s, x + 0.38, y0 + 0.28, cw - 0.7, 0.3, [[("RECOMMENDED", 11, ACCENT, True)]])
    txt(s, x + 0.38, y0 + 0.6, cw - 0.7, 0.5, [[(t, 17, tcol, True)]])
    txt(s, x + 0.38, y0 + 1.05, cw - 0.7, 0.6, [[(price, 24, ACCENT, True)]])
    txt(s, x + 0.38, y0 + 1.75, cw - 0.7, 1.1, [[(d, 12.5, scol, False)]], line_spacing=1.12)
txt(s, 0.9, 6.2, 11.5, 0.8, [
    [("An outside firm would charge this group ", 14, INK, False),
     ("$30,000 to $75,000 or more", 14, ACCENT, True),
     (" to build the same system. You", 14, INK, False)],
    [("already know their floors, so you are the better fit and still a bargain.", 14, INK, False)],
], line_spacing=1.15)

# ----------------------------------------------------------------------------
# Slide 9: First 30 days ramp
# ----------------------------------------------------------------------------
s = slide()
kicker(s, "The first 30 days", "9")
headline(s, ["Your first month, one step at a time."])
weeks = [
    ("Week 1", "Set up Claude.\nFinish one sample document."),
    ("Week 2", "Finish two or three more.\nShape your pitch."),
    ("Week 3", "Have the friendly\nconversation with the group."),
    ("Week 4", "Land the arrangement.\nStart the opening playbook."),
]
baseY = 6.1
heights = [1.4, 2.0, 2.6, 3.2]
cw, gap, x0 = 2.7, 0.45, 1.0
for i, (t, d) in enumerate(weeks):
    x = x0 + i * (cw + gap)
    h = heights[i]
    y = baseY - h
    col = INK if i == 3 else RGBColor(0x2A, 0x4A, 0x63) if i == 2 else SAND if i == 0 else RGBColor(0xD8, 0xC6, 0xAC)
    rect(s, x, y, cw, h, col)
    rect(s, x, y, cw, 0.1, ACCENT)
    tcol = WHITE if i >= 2 else INK
    txt(s, x + 0.25, y + 0.2, cw - 0.5, 0.5, [[(t, 17, tcol, True)]])
    txt(s, x + 0.25, y + 0.7, cw - 0.5, h - 0.8,
        [[(line, 12, tcol, False)] for line in d.split("\n")], line_spacing=1.1, space_after=1)
rect(s, x0, baseY, (cw + gap) * 3 + cw, 0.05, SLATE)
txt(s, 0.9, 6.45, 11.5, 0.5, [[("Progress is measured in documents finished, not hours. Rest days are built in.", 14, SLATE, False)]])

# ----------------------------------------------------------------------------
# Slide 10: Built around your health
# ----------------------------------------------------------------------------
s = slide()
kicker(s, "On your terms", "10")
headline(s, ["Built around your life,", "not the other way around."])
pts = [
    "Work in short bursts. Twenty good minutes moves a document forward.",
    "Stack the easy days. Get ahead when you feel strong, rest when you do not.",
    "Rest the hard days with no guilt. The work waits. It does not expire.",
    "On a retainer, a slow week costs you nothing. You are paid for finished work, not hours.",
]
y = 3.3
for p in pts:
    rect(s, 0.95, y + 0.06, 0.22, 0.22, ACCENT)
    txt(s, 1.45, y, 11.0, 0.6, [[(p, 16.5, INK, False)]], line_spacing=1.08)
    y += 0.78
txt(s, 0.9, 6.7, 11.5, 0.5, [[("You built businesses that ran while you slept. This one runs around your health.", 14, SLATE, False)]])

# ----------------------------------------------------------------------------
# Slide 11: Closing
# ----------------------------------------------------------------------------
s = slide(DEEP)
rect(s, 0, 0, SW, SH, DEEP)
rect(s, 0.9, 2.5, 0.9, 0.16, ACCENT)
txt(s, 0.9, 2.85, 11.5, 2.0, [
    [("You already built this.", 46, WHITE, True)],
    [("Now you sell it.", 46, ACCENT, True)],
], line_spacing=1.03, space_after=2)
txt(s, 0.9, 5.1, 11.3, 1.0, [
    [("Your first move is simple: set up Claude, then read the Start Here playbook.", 18, WHITE, False)],
    [("One finished document at a time, Mark. You have got this.", 18, SLATE, False)],
], line_spacing=1.2)

prs.save("The-Vision.pptx")
print("Saved The-Vision.pptx with", len(prs.slides._sldIdLst), "slides")
